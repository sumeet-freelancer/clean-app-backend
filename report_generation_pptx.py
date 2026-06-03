from datetime import date, datetime, timedelta, timezone
from pathlib import Path
import re

from PIL import Image

JST = timezone(timedelta(hours=9))

SECTION_SLIDE_INDEXES = {
    "parking": 1,
    "bicycle": 2,
    "entrance": 3,
    "floor": 4,
    "stairs": 5,
    "mailbox_garbage_bulletin": 6,
}


def safe_filename(value: str, fallback: str):
    normalized = re.sub(r"[^\w\-]+", "_", str(value or "").strip(), flags=re.ASCII).strip("_")
    return normalized or fallback


def parse_target_month(target_month: str):
    return datetime.strptime(f"{target_month}-01", "%Y-%m-%d").date()


def coerce_date(value):
    if isinstance(value, datetime):
        return value.astimezone(JST).date()
    if isinstance(value, date):
        return value
    return None


def section_map(section_entries):
    return {entry["id"]: entry for entry in section_entries}


def all_work_dates(section_entries):
    dates = []
    for entry in section_entries:
        for value in entry.get("work_dates", []):
            clean = coerce_date(value)
            if clean:
                dates.append(clean)
    return dates


def remove_shape(shape):
    parent = shape._element.getparent()
    if parent is not None:
        parent.remove(shape._element)


def picture_slots(slide):
    slots = []
    for shape in slide.shapes:
        name = str(getattr(shape, "name", "") or "")
        if "図プレースホルダー" in name or "コンテンツ プレースホルダー" in name:
            slots.append(shape)
    return sorted(slots, key=lambda shape: (shape.top, shape.left))


def add_cropped_picture(slide, image_path, slot):
    left, top, width, height = slot.left, slot.top, slot.width, slot.height
    remove_shape(slot)

    picture = slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    with Image.open(image_path) as image:
        image_width, image_height = image.size

    if image_width <= 0 or image_height <= 0 or width <= 0 or height <= 0:
        return picture

    image_ratio = image_width / image_height
    frame_ratio = width / height
    if image_ratio > frame_ratio:
        crop = max(0, (1 - frame_ratio / image_ratio) / 2)
        picture.crop_left = crop
        picture.crop_right = crop
    elif image_ratio < frame_ratio:
        crop = max(0, (1 - image_ratio / frame_ratio) / 2)
        picture.crop_top = crop
        picture.crop_bottom = crop
    return picture


def replace_text(slide, replacements):
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text
                for old, new in replacements.items():
                    if old in text:
                        text = text.replace(old, new)
                run.text = text


def set_cover_dates(slide, shooting_date, created_date):
    values = [
        shooting_date.year,
        shooting_date.month,
        shooting_date.day,
        created_date.year,
        created_date.month,
        created_date.day,
    ]

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        digit_runs = []
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if str(run.text or "").strip().isdigit():
                    digit_runs.append(run)

        if len(digit_runs) >= 6:
            for run, value in zip(digit_runs[:6], values):
                run.text = str(value)
            return


def fill_photo_slots(slide, photos):
    slots = picture_slots(slide)
    for slot, photo_path in zip(slots, photos):
        add_cropped_picture(slide, Path(photo_path), slot)
    for slot in slots[len(photos):]:
        remove_shape(slot)


def delete_slides(presentation, indexes):
    slide_ids = presentation.slides._sldIdLst
    for index in sorted(indexes, reverse=True):
        if index < 0 or index >= len(slide_ids):
            continue
        rel_id = slide_ids[index].rId
        presentation.part.drop_rel(rel_id)
        del slide_ids[index]


def validate_required_sections(section_entries):
    missing = [
        entry.get("label") or entry["id"]
        for entry in section_entries
        if entry.get("required") and not entry.get("photos")
    ]
    if missing:
        raise RuntimeError(f"missing required report sections: {', '.join(missing)}")


def generate_visual_inspection_report(
    *,
    property_name: str,
    target_month: str,
    template_path: Path,
    output_dir: Path,
    section_entries,
):
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for pptx-visual-inspection reports") from exc

    if not template_path.exists():
        raise FileNotFoundError(f"template not found: {template_path}")

    validate_required_sections(section_entries)
    output_dir.mkdir(parents=True, exist_ok=True)

    presentation = Presentation(str(template_path))
    entries = section_map(section_entries)
    target_date = parse_target_month(target_month)
    work_dates = all_work_dates(section_entries)
    shooting_date = min(work_dates) if work_dates else target_date
    created_date = datetime.now(JST).date()

    cover_slide = presentation.slides[0]
    replace_text(cover_slide, {"サンプル": property_name or "物件"})
    set_cover_dates(cover_slide, shooting_date, created_date)

    exterior_photos = entries.get("exterior", {}).get("photos", [])
    if exterior_photos:
        fill_photo_slots(cover_slide, exterior_photos[:1])

    slides_to_delete = set()
    for section_id, slide_index in SECTION_SLIDE_INDEXES.items():
        entry = entries.get(section_id, {})
        photos = entry.get("photos", [])
        if slide_index >= len(presentation.slides):
            if photos:
                raise RuntimeError(f"template missing slide for section: {section_id}")
            continue

        if photos:
            fill_photo_slots(presentation.slides[slide_index], photos)
        else:
            slides_to_delete.add(slide_index)

    delete_slides(presentation, slides_to_delete)

    property_token = safe_filename(property_name, "visual_inspection")
    output_name = f"{target_month}_{property_token}_visual_inspection_report.pptx"
    output_path = output_dir / output_name
    presentation.save(output_path)
    return output_path
